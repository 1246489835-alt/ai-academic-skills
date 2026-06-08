"""
CNS-PPT-排版 通用模板脚本
========================
Nature/Cell/Science 投稿级 Figure PPT 自动排版。
使用方法：修改下方 CONFIG 区的路径，运行即可。

命名规则: 【letter-number】描述.tif
         【letter】描述.tif (单张panel)
"""

import os
import re
import sys
from pathlib import Path
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image


# ==================== CONFIG ====================
# 修改这里的路径即可复用

INPUT_DIR = r""          # 图片文件夹（如果有【1】主图/【2】附图子文件夹会自动识别）
OUTPUT_DIR = r""         # 输出目录
OUTPUT_NAME = ""         # 输出文件名，如 "2026-05-27-R3-Figure-Layout.pptx"

# 页面参数（一般不用改）
SLIDE_W = Inches(7.5)    # 纵向A4
SLIDE_H = Inches(10.0)
M_LEFT = Cm(0.7)
M_RIGHT = Cm(0.7)
M_TOP = Cm(0.7)
M_BOT = Cm(0.4)
PANEL_VGAP = Cm(0.45)   # panel间纵向间距
IMG_HGAP = Cm(0.2)      # panel内图片横向间距
LABEL_SIZE = Pt(12)      # panel标签字号
LABEL_FONT = "Arial"
LABEL_H = Cm(0.45)      # 标签占用高度


# ==================== CORE ====================

def parse_filename(filename):
    """解析【letter-number】或【letter】命名"""
    match = re.match(r'[【\[]([a-zA-Z])-?(\d*)[】\]](.+)', filename)
    if match:
        panel = match.group(1).lower()
        order = int(match.group(2)) if match.group(2) else 1
        desc = match.group(3)
        return panel, order, desc
    return None, None, None


def scan_folder(folder_path):
    """扫描文件夹，按panel分组返回"""
    if not os.path.isdir(folder_path):
        return {}
    panels = defaultdict(list)
    exts = {'.tif', '.tiff', '.png', '.jpg', '.jpeg', '.bmp'}
    for f in Path(folder_path).iterdir():
        if f.suffix.lower() in exts:
            panel, order, desc = parse_filename(f.name)
            if panel is not None:
                with Image.open(str(f)) as img:
                    w, h = img.size
                panels[panel].append({
                    'path': str(f), 'order': order,
                    'w_px': w, 'h_px': h, 'ar': w / h
                })
    for p in panels:
        panels[p].sort(key=lambda x: x['order'])
    return dict(sorted(panels.items()))


def detect_subfolders(input_dir):
    """检测是否有【1】主图 / 【2】附图子文件夹"""
    main_dir = None
    supp_dir = None
    for d in Path(input_dir).iterdir():
        if d.is_dir():
            name = d.name
            if '主图' in name or name.startswith('【1】'):
                main_dir = str(d)
            elif '附图' in name or name.startswith('【2】'):
                supp_dir = str(d)
    return main_dir, supp_dir


def fit(ar, max_w, max_h):
    """保持宽高比缩放"""
    if ar * max_h <= max_w:
        h = max_h; w = int(h * ar)
    else:
        w = max_w; h = int(w / ar)
    return w, h


def ideal_height(panel_key, imgs, avail_w, hgap):
    """计算panel以全宽排列时的理想高度"""
    n = len(imgs)

    # 检测特殊组合：首图AR明显不同于其余（如UMAP+barplots）
    if n >= 4 and imgs[0]['ar'] < 1.5 and all(i['ar'] > 1.5 for i in imgs[1:]):
        left_w = avail_w * 0.40
        right_w = avail_w * 0.58
        h_left = left_w / imgs[0]['ar']
        cell_w = (right_w - hgap) / 2
        avg_ar = sum(i['ar'] for i in imgs[1:]) / (n - 1)
        rows = ((n - 1) + 1) // 2
        h_right = rows * (cell_w / avg_ar) + hgap * (rows - 1)
        return max(h_left, h_right)

    # 检测方形+宽幅组合
    if n == 2 and any(abs(i['ar'] - 1.0) < 0.15 for i in imgs):
        sq = min(imgs, key=lambda x: abs(x['ar'] - 1.0))
        sq_w = avail_w * 0.38
        return sq_w / sq['ar']

    # 默认：全部并排
    if n == 1:
        return avail_w / imgs[0]['ar']
    elif n <= 3:
        slot_w = (avail_w - hgap * (n - 1)) / n
        min_ar = min(i['ar'] for i in imgs)
        return slot_w / min_ar
    else:
        cols = min(n, 3)
        rows = (n + cols - 1) // cols
        slot_w = (avail_w - hgap * (cols - 1)) / cols
        min_ar = min(i['ar'] for i in imgs)
        return rows * (slot_w / min_ar) + hgap * (rows - 1)


def add_label(slide, letter, x, y):
    """Nature风格bold lowercase panel标签"""
    box = slide.shapes.add_textbox(x, y, int(Cm(0.8)), int(Cm(0.5)))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = letter
    run.font.name = LABEL_FONT
    run.font.size = LABEL_SIZE
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_figure_slide(prs, panels):
    """为一组panels生成一页slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)

    x0 = int(M_LEFT)
    y0 = int(M_TOP)
    W = int(SLIDE_W) - int(M_LEFT) - int(M_RIGHT)
    H = int(SLIDE_H) - int(M_TOP) - int(M_BOT)
    vgap = int(PANEL_VGAP)
    hgap = int(IMG_HGAP)
    lbl_h = int(LABEL_H)

    panel_keys = list(panels.keys())
    n_panels = len(panel_keys)
    if n_panels == 0:
        return slide

    # 计算理想高度并归一化
    ideals = {}
    for key in panel_keys:
        ideals[key] = ideal_height(key, panels[key], W, hgap)

    overhead = lbl_h * n_panels + vgap * (n_panels - 1)
    content_h = H - overhead
    total_ideal = sum(ideals.values())
    scale = content_h / total_ideal if total_ideal > 0 else 1.0
    row_heights = {k: int(v * scale) for k, v in ideals.items()}

    # 逐panel放置
    cur_y = y0
    for key in panel_keys:
        imgs = panels[key]
        row_h = row_heights[key]
        n = len(imgs)

        add_label(slide, key, x0, cur_y)
        img_y = cur_y + lbl_h

        # --- 布局策略 ---

        # 策略A: 首图+网格 (如UMAP左 + barplots右)
        if n >= 4 and imgs[0]['ar'] < 1.5 and all(i['ar'] > 1.5 for i in imgs[1:]):
            left_w = int(W * 0.40)
            gap_mid = int(Cm(0.3))
            right_w = W - left_w - gap_mid

            iw, ih = fit(imgs[0]['ar'], left_w, row_h)
            ix = x0 + (left_w - iw) // 2
            iy = img_y + (row_h - ih) // 2
            slide.shapes.add_picture(imgs[0]['path'], ix, iy, width=iw, height=ih)

            rest = imgs[1:]
            cols = 2
            rrows = (len(rest) + cols - 1) // cols
            cell_w = (right_w - hgap * (cols - 1)) // cols
            cell_h = (row_h - hgap * (rrows - 1)) // rrows
            rx0 = x0 + left_w + gap_mid
            for idx, img in enumerate(rest):
                r = idx // cols
                c = idx % cols
                iw, ih = fit(img['ar'], cell_w, cell_h)
                ix = rx0 + c * (cell_w + hgap) + (cell_w - iw) // 2
                iy2 = img_y + r * (cell_h + hgap) + (cell_h - ih) // 2
                slide.shapes.add_picture(img['path'], ix, iy2, width=iw, height=ih)

        # 策略B: 方形+宽幅组合
        elif n == 2 and any(abs(i['ar'] - 1.0) < 0.15 for i in imgs):
            sq = min(imgs, key=lambda x: abs(x['ar'] - 1.0))
            wide = [i for i in imgs if i is not sq][0]

            sq_area_w = int(W * 0.38)
            gap_mid = int(Cm(0.3))
            wide_area_w = W - sq_area_w - gap_mid

            iw, ih = fit(sq['ar'], sq_area_w, row_h)
            ix = x0 + (sq_area_w - iw) // 2
            iy = img_y + (row_h - ih) // 2
            slide.shapes.add_picture(sq['path'], ix, iy, width=iw, height=ih)

            iw, ih = fit(wide['ar'], wide_area_w, row_h)
            ix = x0 + sq_area_w + gap_mid + (wide_area_w - iw) // 2
            iy = img_y + (row_h - ih) // 2
            slide.shapes.add_picture(wide['path'], ix, iy, width=iw, height=ih)

        # 策略C: 单张居中
        elif n == 1:
            iw, ih = fit(imgs[0]['ar'], W, row_h)
            ix = x0 + (W - iw) // 2
            iy = img_y + (row_h - ih) // 2
            slide.shapes.add_picture(imgs[0]['path'], ix, iy, width=iw, height=ih)

        # 策略D: 2-3张并排
        elif n <= 3:
            slot_w = (W - hgap * (n - 1)) // n
            for i, img in enumerate(imgs):
                iw, ih = fit(img['ar'], slot_w, row_h)
                ix = x0 + i * (slot_w + hgap) + (slot_w - iw) // 2
                iy = img_y + (row_h - ih) // 2
                slide.shapes.add_picture(img['path'], ix, iy, width=iw, height=ih)

        # 策略E: 4+张网格
        else:
            cols = min(n, 3)
            rrows = (n + cols - 1) // cols
            cell_w = (W - hgap * (cols - 1)) // cols
            cell_h = (row_h - hgap * (rrows - 1)) // rrows
            for idx, img in enumerate(imgs):
                r = idx // cols
                c = idx % cols
                iw, ih = fit(img['ar'], cell_w, cell_h)
                ix = x0 + c * (cell_w + hgap) + (cell_w - iw) // 2
                iy = img_y + r * (cell_h + hgap) + (cell_h - ih) // 2
                slide.shapes.add_picture(img['path'], ix, iy, width=iw, height=ih)

        cur_y = img_y + row_h + vgap

    return slide


# ==================== MAIN ====================

def main():
    print("=" * 50)
    print("CNS-PPT Nature-Grade Figure Layout")
    print("=" * 50)

    # 检测子文件夹结构
    main_dir, supp_dir = detect_subfolders(INPUT_DIR)

    if main_dir or supp_dir:
        # 方式二：主图+附图分文件夹
        main_panels = scan_folder(main_dir) if main_dir else {}
        supp_panels = scan_folder(supp_dir) if supp_dir else {}
    else:
        # 方式一：单文件夹
        main_panels = scan_folder(INPUT_DIR)
        supp_panels = {}

    # 打印扫描结果
    if main_panels:
        print("\n[Main Figure]")
        for k, imgs in main_panels.items():
            ars = ", ".join(["%.2f" % i['ar'] for i in imgs])
            print(f"  Panel {k}: {len(imgs)} imgs  AR=[{ars}]")

    if supp_panels:
        print("\n[Supplementary Figure]")
        for k, imgs in supp_panels.items():
            ars = ", ".join(["%.2f" % i['ar'] for i in imgs])
            print(f"  Panel {k}: {len(imgs)} imgs  AR=[{ars}]")

    if not main_panels and not supp_panels:
        print("[ERROR] No images found matching naming pattern.")
        sys.exit(1)

    # 生成PPTX
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    if main_panels:
        add_figure_slide(prs, main_panels)
    if supp_panels:
        add_figure_slide(prs, supp_panels)

    out = os.path.join(OUTPUT_DIR, OUTPUT_NAME)
    prs.save(out)

    n_slides = len(prs.slides)
    print(f"\n-> {out}")
    print(f"   {n_slides} slide(s) | Portrait 7.5x10in | Nature style")
    print("=" * 50)


if __name__ == "__main__":
    main()
